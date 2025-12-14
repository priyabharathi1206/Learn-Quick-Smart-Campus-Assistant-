import pdfplumber
from pptx import Presentation
from docx import Document
from sentence_transformers import SentenceTransformer
import numpy as np
import faiss
import re
from groq import Groq
import os
from dotenv import load_dotenv

load_dotenv()

# Groq Client
client = Groq(api_key=os.getenv("GROQ_API_KEY"))

# Text Extraction 
def extract_text_from_pdf(file_path):
    text = ""
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            t = page.extract_text() or ""
            text += t + "\n"
    return text

def extract_text_from_docx(file_path):
    doc = Document(file_path)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

#  Cleaning & Chunking 

def clean_text(text):
    return re.sub(r'\s+', ' ', text)

def split_into_chunks(text, chunk_size=500, overlap=50):
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start += chunk_size - overlap
    return chunks

#  Embeddings + FAISS 

embed_model = SentenceTransformer("all-MiniLM-L6-v2")

def create_vector_store(chunks):
    embeddings = embed_model.encode(chunks)

    if embeddings is None or len(embeddings) == 0:
        raise ValueError("Embedding failed — chunks are empty.")

    embedding_dim = embeddings.shape[1]

    index = faiss.IndexFlatL2(embedding_dim)
    index.add(np.array(embeddings))

    return embeddings, index

#  RAG Answering 
def get_answer(question, context):
    prompt = f"""
You are a Smart Campus Assistant.

Use ONLY the following study material to answer:

CONTEXT:
{context}

QUESTION:
{question}

Give a clear answer.
"""
    response = client.chat.completions.create(
        model="llama-3.1-8b-instant",
        messages=[{"role": "user", "content": prompt}],
        max_tokens=300
    )
    return response.choices[0].message.content


def ask_question(question, chunks, index, top_k=5):
    q_embed = embed_model.encode([question])
    D, I = index.search(np.array(q_embed), top_k)

    retrieved = "\n\n".join([chunks[i] for i in I[0]])
    return get_answer(question, retrieved)

#  MCQ Generation 

def generate_mcqs(context, num_questions=5):
    prompt = f"""
Generate {num_questions} multiple-choice questions from the study material below.

Format:
Q1: <question>
A) option
B) option
C) option
D) option
Correct Answer: <A/B/C/D>
Explanation: <why>

STUDY MATERIAL:
{context}
"""

    response = client.chat.completions.create(
        model="llama-3.1-8b-instant",
        messages=[{"role": "user", "content": prompt}],
        max_tokens=800
    )

    return response.choices[0].message.content



def parse_mcq_text(text):
    mcqs = []
    blocks = text.split("Q")[1:]  # split per question

    for block in blocks:
        lines = block.strip().split("\n")
        q_line = lines[0][2:]   # remove "1:" or "2:"
        options = []
        correct = None

        for line in lines:
            if line.startswith(("A)", "B)", "C)", "D)")):
                options.append(line[2:].strip())

            if "Correct Answer" in line:
                correct_letter = line.split(":")[1].strip()
                correct = "ABCD".index(correct_letter)

        mcqs.append({
            "question": q_line,
            "options": options,
            "answer": correct
        })

    return mcqs




def get_mcqs(chunks, index, num_q=5, top_k=5):
    question = "Generate exam questions"
    q_embed = embed_model.encode([question])
    D, I = index.search(np.array(q_embed), top_k)

    retrieved = "\n".join([chunks[i] for i in I[0]])

    raw = generate_mcqs(retrieved, num_q)

    return parse_mcq_text(raw)

def check_answer(mcq_text, q_no, user_answer):
    # Match Q1, Q1., or Q1: and allow optional spaces before Correct Answer
    pattern = rf"Q{q_no}[:.]?.*?Correct Answer[:\s]*([A-D])"
    match = re.search(pattern, mcq_text, re.DOTALL | re.IGNORECASE)

    if not match:
        return "Correct answer not found."

    correct = match.group(1).strip().upper()
    user_answer = user_answer.upper()

    if user_answer == correct:
        return "Correct! "
    return f"Wrong  — Correct answer is {correct}"

#  Summary Function 

def summarize_text(text):
    prompt = f"""
Summarize the following clearly and concisely:

TEXT:
{text}

SUMMARY:
"""

    response = client.chat.completions.create(
        model="llama-3.1-8b-instant",
        messages=[{"role": "user", "content": prompt}],
        max_tokens=300
    )
    return response.choices[0].message.content

def extract_topics_and_keywords(text, max_topics=10):
    """
    Extract main topics and their keywords from study material.
    Returns a list of dictionaries: [{"topic": "X", "keywords": ["a","b","c"]}, ...]
    """
    prompt = f"""
You are a Smart Campus Assistant.

Analyze the following study material and list the main topics and their subtopics with relevant keywords.
- Provide at most {max_topics} topics.
- Each topic should have a list of 3-7 keywords that are important for understanding the topic.
- Output in JSON format ONLY like this:

[
  {{
    "topic": "Topic Name",
    "keywords": ["keyword1", "keyword2", "keyword3"]
  }},
  ...
]

STUDY MATERIAL:
{text}
"""
    response = client.chat.completions.create(
        model="llama-3.1-8b-instant",
        messages=[{"role": "user", "content": prompt}],
        max_tokens=600
    )

    result_text = response.choices[0].message.content.strip()

    if not result_text:
        print("Groq API returned empty response")
        return [{"topic": "No topics found", "keywords": []}]

    # Sometimes API returns extra text before/after JSON, try to extract JSON part
    try:
        import json
        # Find first [ and last ] to extract JSON
        start = result_text.find("[")
        end = result_text.rfind("]") + 1
        json_text = result_text[start:end]
        topics = json.loads(json_text)
        if not isinstance(topics, list):
            raise ValueError("Parsed result is not a list")
        return topics
    except Exception as e:
        print("Error parsing topics:", e, "API output:", result_text)
        return [{"topic": "Parsing failed", "keywords": []}]


def extract_topics_hierarchy(text, max_topics=10):
    prompt = f"""
You are a Smart Campus Assistant.

Analyze the following study material and generate a hierarchical structure:

- Maximum {max_topics} main topics
- Each topic can have 2-5 subtopics
- Include important keywords under each subtopic
- Output **JSON only** like this:

[
  {{
    "topic": "Main Topic",
    "subtopics": [
        {{"name": "Subtopic 1", "keywords": ["k1","k2","k3"]}},
        {{"name": "Subtopic 2", "keywords": ["k4","k5"]}}
    ]
  }}
]

STUDY MATERIAL:
{text}
"""
    response = client.chat.completions.create(
        model="llama-3.1-8b-instant",
        messages=[{"role": "user", "content": prompt}],
        max_tokens=800
    )

    result_text = response.choices[0].message.content.strip()

    if not result_text:
        print("Groq API returned empty response")
        return [{"topic": "No topics found", "subtopics": []}]

    try:
        import json, re
        # Extract first JSON array in the text
        match = re.search(r'\[.*\]', result_text, re.DOTALL)
        if not match:
            raise ValueError("No JSON array found in API response")
        json_text = match.group()
        topics = json.loads(json_text)
        return topics
    except Exception as e:
        print("Error parsing hierarchical topics:", e, "API output:", result_text)
        return [{"topic": "Parsing failed", "subtopics": []}]
